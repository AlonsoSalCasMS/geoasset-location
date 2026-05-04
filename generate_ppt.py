"""
Generates GeoAssets_Intelligence_TFM.pptx
Visual style based on the Management Solutions TFM example.
"""

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Slide dimensions (matching example PPT) ─────────────────────
SW = 33.87
SH = 19.05

# ── Color palette (from example) ────────────────────────────────
NAVY      = RGBColor(0x00, 0x1A, 0x48)
BLUE      = RGBColor(0x00, 0x25, 0x69)
GOLD      = RGBColor(0xCA, 0xB4, 0x51)
DARK_GOLD = RGBColor(0x99, 0x7D, 0x00)
DARK2     = RGBColor(0x2B, 0x35, 0x49)
GREY      = RGBColor(0xF2, 0xF2, 0xF2)
CREAM     = RGBColor(0xFB, 0xF9, 0xF3)
GEO_BLUE  = RGBColor(0x21, 0x96, 0xF3)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x1A, 0x1A, 0x1A)
MID_GREY  = RGBColor(0x75, 0x75, 0x75)
RED_C     = RGBColor(0xF4, 0x43, 0x36)
ORG_C     = RGBColor(0xFF, 0x98, 0x00)
GRN_C     = RGBColor(0x4C, 0xAF, 0x50)
TEAL_C    = RGBColor(0x00, 0x96, 0x88)
PURPLE_C  = RGBColor(0x7B, 0x1F, 0xA2)

# ── Layout constants (cm) ───────────────────────────────────────
HDR_H   = 3.1    # header bar height
ML      = 1.5    # left margin
MR      = 1.5    # right margin
CW      = SW - ML - MR   # content width = 30.87
CY      = 5.0    # content start y
CH      = SH - CY - 1.0  # content height available
FOOTER_Y = SH - 0.8


# ════════════════════════════════════════════════════════════════
# PRIMITIVES
# ════════════════════════════════════════════════════════════════

def _rect(slide, l, t, w, h, fill=None, border=None, border_w_pt=0, rounding=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounding else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Cm(l), Cm(t), Cm(w), Cm(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border and border_w_pt > 0:
        s.line.color.rgb = border
        s.line.width = Pt(border_w_pt)
    else:
        s.line.fill.background()
    if rounding:
        try: s.adjustments[0] = 0.08
        except: pass
    return s


def _txt(slide, l, t, w, h, text, size, bold=False, italic=False,
         color=BLACK, align=PP_ALIGN.LEFT, wrap=True, spacing_after=0):
    box = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing_after:
        p.space_after = Pt(spacing_after)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def _shape_txt(shape, lines, size, bold=False, italic=False,
               color=WHITE, align=PP_ALIGN.LEFT, v_anchor=None):
    """Set multi-line text inside a shape."""
    from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
    tf = shape.text_frame
    tf.word_wrap = True
    if v_anchor:
        tf.vertical_anchor = v_anchor
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if isinstance(line, tuple):
            # (text, bold_override)
            run = p.add_run()
            run.text = line[0]
            run.font.size = Pt(size)
            run.font.bold = line[1] if len(line) > 1 else bold
            run.font.color.rgb = line[2] if len(line) > 2 else color
            run.font.italic = italic
        else:
            run = p.add_run()
            run.text = line
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color


def _hline(slide, l, t, w, color=GOLD, height_pt=1.5):
    s = _rect(slide, l, t, w, height_pt / 28.35, fill=color)
    return s


# ════════════════════════════════════════════════════════════════
# COMPOSITE COMPONENTS
# ════════════════════════════════════════════════════════════════

def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def add_bg(slide, color=WHITE):
    bg = _rect(slide, 0, 0, SW, SH, fill=color)
    # send to back via XML
    sp = bg._element
    slide.shapes._spTree.remove(sp)
    slide.shapes._spTree.insert(2, sp)


def add_header(slide, sec_num, sec_name, slide_title):
    """Navy header bar with gold accent and section/title labels."""
    # Navy bar
    _rect(slide, 0, 0, SW, HDR_H, fill=NAVY)
    # Gold left accent
    _rect(slide, 0, 0, 0.45, HDR_H, fill=GOLD)
    # Section badge background
    badge_w = len(f"SECCIÓN {sec_num}  ·  {sec_name}") * 0.18 + 1.5
    badge_w = min(badge_w, 22.0)
    s_badge = _rect(slide, 0.75, 0.45, badge_w, 0.75, fill=DARK_GOLD, rounding=True)
    _shape_txt(s_badge, [f"SECCIÓN {sec_num}  ·  {sec_name}"], 8.5, bold=True,
               color=WHITE, align=PP_ALIGN.LEFT)
    # Slide title in header
    _txt(slide, 0.75, 1.35, SW - 1.5, 1.5, slide_title,
         size=17, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def add_subtitle_bar(slide, text):
    """Context sentence below header."""
    s = _rect(slide, 0, HDR_H, SW, 1.5, fill=CREAM)
    _shape_txt(s, [text], 10, bold=True, color=BLUE, align=PP_ALIGN.LEFT)
    tf = s.text_frame
    tf.paragraphs[0].runs[0].font.italic = True
    # inner padding
    from pptx.oxml.ns import qn
    txBody = s.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('lIns', str(int(Cm(0.4))))
        bodyPr.set('tIns', str(int(Cm(0.3))))


def add_footer_line(slide, slide_num_text=""):
    _hline(slide, 0, FOOTER_Y, SW, color=NAVY, height_pt=1.5)
    if slide_num_text:
        _txt(slide, SW - 3.5, FOOTER_Y + 0.05, 3.2, 0.65,
             slide_num_text, 8, color=WHITE, align=PP_ALIGN.RIGHT)
    _txt(slide, 0.6, FOOTER_Y + 0.05, 20, 0.65,
         "GeoAssets Intelligence  ·  TFM Máster en Inteligencia Artificial  ·  2025-2026",
         8, color=WHITE)


def card(slide, l, t, w, h, header_text, body_lines,
         header_color=NAVY, body_color=GREY, txt_color=WHITE,
         body_txt_color=BLACK, header_size=10, body_size=9.5, rounding=True):
    """A card with colored header and lighter body."""
    hdr_h = 0.85
    # Header
    s_hdr = _rect(slide, l, t, w, hdr_h, fill=header_color, rounding=rounding)
    _shape_txt(s_hdr, [header_text], header_size, bold=True, color=txt_color,
               align=PP_ALIGN.CENTER)
    # Body
    body_h = h - hdr_h
    s_body = _rect(slide, l, t + hdr_h, w, body_h, fill=body_color,
                   border=header_color, border_w_pt=0.5, rounding=False)
    _shape_txt(s_body, body_lines, body_size, color=body_txt_color, align=PP_ALIGN.LEFT)
    tf = s_body.text_frame
    from pptx.oxml.ns import qn
    txBody = tf._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('lIns', str(int(Cm(0.25))))
        bodyPr.set('tIns', str(int(Cm(0.2))))
    return s_hdr, s_body


def two_col(slide, y, left_hdr, left_items, right_hdr, right_items,
            left_color=NAVY, right_color=DARK_GOLD, h=7.5):
    """Two-column layout with header row and bullet ovals."""
    col_w = (CW - 0.5) / 2
    lx = ML
    rx = ML + col_w + 0.5

    # Left column
    s_lhdr = _rect(slide, lx, y, col_w, 0.9, fill=left_color, rounding=False)
    _shape_txt(s_lhdr, [left_hdr], 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    s_lbody = _rect(slide, lx, y + 0.9, col_w, h - 0.9, fill=GREY, border=left_color, border_w_pt=0.8)
    for i, item in enumerate(left_items):
        iy = y + 0.9 + 0.25 + i * 1.25
        if iy + 0.8 > y + h: break
        bullet = _rect(slide, lx + 0.25, iy, 0.65, 0.65, fill=left_color, rounding=True)
        _shape_txt(bullet, [str(i + 1)], 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(slide, lx + 1.05, iy - 0.05, col_w - 1.3, 0.8,
             item, 9.5, color=BLACK)

    # Right column
    s_rhdr = _rect(slide, rx, y, col_w, 0.9, fill=right_color, rounding=False)
    _shape_txt(s_rhdr, [right_hdr], 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    s_rbody = _rect(slide, rx, y + 0.9, col_w, h - 0.9, fill=CREAM, border=right_color, border_w_pt=0.8)
    for i, item in enumerate(right_items):
        iy = y + 0.9 + 0.25 + i * 1.25
        if iy + 0.8 > y + h: break
        bullet = _rect(slide, rx + 0.25, iy, 0.65, 0.65, fill=right_color, rounding=True)
        _shape_txt(bullet, [str(i + 1)], 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(slide, rx + 1.05, iy - 0.05, col_w - 1.3, 0.8,
             item, 9.5, color=BLACK)


def pipeline_step(slide, x, y, w, h, num, title, desc, color=NAVY):
    """Single pipeline step block with number badge."""
    badge_r = 0.55
    # Step box
    s = _rect(slide, x, y, w, h, fill=color, rounding=True)
    _shape_txt(s, [title, "", desc], 9.5, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    tf = s.text_frame
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.size = Pt(10)
    # Number badge (top-center)
    badge_x = x + w / 2 - badge_r / 2
    b = _rect(slide, badge_x, y - badge_r * 0.5, badge_r, badge_r, fill=GOLD, rounding=True)
    _shape_txt(b, [num], 10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    return s


def arrow_right(slide, x, y, color=GOLD):
    """Small right-pointing arrow between pipeline steps."""
    s = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        Cm(x), Cm(y), Cm(0.8), Cm(0.5)
    )
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


# ════════════════════════════════════════════════════════════════
# SPECIAL SLIDE TYPES
# ════════════════════════════════════════════════════════════════

def section_divider(prs, sec_num, sec_name, description=""):
    """Full navy slide for section transitions."""
    slide = add_slide(prs)
    add_bg(slide, NAVY)

    # Gold decorative horizontal lines
    _hline(slide, 0, 2.8, SW, GOLD, 3)
    _hline(slide, 0, 3.25, SW, DARK_GOLD, 1.5)
    _hline(slide, 0, SH - 3.2, SW, GOLD, 3)
    _hline(slide, 0, SH - 3.65, SW, DARK_GOLD, 1.5)

    # Large section number
    _txt(slide, 1.5, 4.5, 10, 4.5, sec_num,
         size=72, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    # Section title
    _txt(slide, 1.5, 8.5, SW - 3, 3.0, sec_name,
         size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Description
    if description:
        _txt(slide, 1.5, 11.5, SW - 3, 2.5, description,
             size=13, italic=True, color=GOLD, align=PP_ALIGN.LEFT)

    # Footer
    _txt(slide, SW - 8, SH - 1.8, 7, 1.0,
         "GeoAssets Intelligence · TFM 2025-2026",
         9, italic=True, color=DARK_GOLD, align=PP_ALIGN.RIGHT)

    add_footer_line(slide)
    return slide


def content_slide(prs, sec_num, sec_name, slide_title, subtitle=""):
    slide = add_slide(prs)
    add_header(slide, sec_num, sec_name, slide_title)
    if subtitle:
        add_subtitle_bar(slide, subtitle)
    add_footer_line(slide)
    return slide


# ════════════════════════════════════════════════════════════════
# SLIDE CONTENT BUILDERS
# ════════════════════════════════════════════════════════════════

def slide_cover(prs):
    slide = add_slide(prs)
    add_bg(slide, NAVY)

    # Gold accent rectangle (top-right decorative)
    _rect(slide, SW - 9, 0, 9, 7.5, fill=DARK_GOLD)
    _rect(slide, SW - 9, 0, 0.5, 7.5, fill=GOLD)
    _hline(slide, 0, 7.5, SW, GOLD, 4)
    _hline(slide, 0, 8.1, SW, DARK_GOLD, 2)

    # Logo area placeholder
    _txt(slide, SW - 8.5, 0.5, 7.5, 1.5,
         "MÁSTER EN INTELIGENCIA ARTIFICIAL", 9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    _txt(slide, SW - 8.5, 1.7, 7.5, 3.5,
         "Trabajo\nde Fin\nde Máster", 20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Main title
    _txt(slide, 1.5, 2.0, SW - 12, 3.5,
         "GeoAssets\nIntelligence", 38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle gold line + text
    _hline(slide, 1.5, 5.6, 14, GOLD, 3)
    _txt(slide, 1.5, 6.2, 20, 1.5,
         "Localización automatizada del patrimonio\nproductivo empresarial mediante IA",
         14, italic=True, color=GOLD, align=PP_ALIGN.LEFT)

    # Tags row
    for i, (lbl, clr) in enumerate([
        ("Google Maps API", GEO_BLUE), ("AWS Bedrock", ORG_C), ("CrewAI Agents", TEAL_C)
    ]):
        bx = 1.5 + i * 7.0
        b = _rect(slide, bx, 8.6, 6.5, 0.75, fill=clr, rounding=True)
        _shape_txt(b, [lbl], 9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Author / date
    _txt(slide, 1.5, 10.0, 20, 1.0,
         "Alonso Salgueiro Castelo  ·  Tutor: —  ·  Mayo 2026",
         10, color=MID_GREY, align=PP_ALIGN.LEFT)

    # Footer
    _hline(slide, 0, SH - 0.8, SW, NAVY, 1.5)
    _txt(slide, 0.6, SH - 0.75, 25, 0.65,
         "GeoAssets Intelligence  ·  TFM Máster en Inteligencia Artificial  ·  2025-2026",
         8, color=MID_GREY)


def slide_index(prs):
    slide = add_slide(prs)
    add_bg(slide, DARK2)
    # Dark overlay
    _rect(slide, 0, 0, SW, SH, fill=NAVY)

    # Gold decorative bar
    _rect(slide, 0, 0, 0.5, SH, fill=GOLD)
    _hline(slide, 0.5, 3.2, SW, GOLD, 2)

    _txt(slide, 1.5, 1.0, 20, 2.0, "Índice", 36, bold=True, color=WHITE)

    sections = [
        ("01", "El problema: opacidad del patrimonio empresarial", "Slides 3–7"),
        ("02", "Solución: plataforma GeoAssets Intelligence", "Slides 8–11"),
        ("03", "Arquitectura técnica", "Slides 12–15"),
        ("04", "Pipelines de análisis", "Slides 16–22"),
        ("05", "Demostración y resultados", "Slides 23–26"),
        ("06", "Benchmarking competitivo", "Slides 27–28"),
        ("07", "Propuesta de valor y conclusiones", "Slides 29–34"),
    ]

    for i, (num, name, slides) in enumerate(sections):
        y = 3.6 + i * 1.9
        # Badge
        b = _rect(slide, 1.5, y, 1.5, 0.75, fill=DARK_GOLD, rounding=True)
        _shape_txt(b, [num], 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Name
        _txt(slide, 3.3, y, 22, 0.8, name, 12, bold=True, color=WHITE)
        _txt(slide, 3.3, y + 0.6, 22, 0.65, slides, 9, italic=True, color=GOLD)
        # Thin separator
        if i < len(sections) - 1:
            _hline(slide, 1.5, y + 1.55, CW, DARK_GOLD, 0.8)

    add_footer_line(slide)


def slide_s01_context(prs):
    s = content_slide(prs, "01", "El problema",
                      "Contexto: los activos productivos empresariales",
                      "Las empresas poseen activos distribuidos en múltiples localizaciones — conocerlos es clave para el análisis de riesgo, M&A y ESG.")
    y = CY + 0.2
    # 4 sector cards
    sectors = [
        ("Banca & Seguros", ["Valoración de colaterales", "Riesgo de concentración geográfica", "Due diligence en M&A"]),
        ("Consultoría Estratégica", ["Mapa competitivo de activos", "Inteligencia de mercado", "Expansión geográfica"]),
        ("ESG / Sostenibilidad", ["Exposición ambiental industrial", "Huella de carbono por activo", "Reporte de sostenibilidad"]),
        ("Inversión & Capital", ["Valoración de carteras inmobiliarias", "Screening de targets", "Monitorización de portfolio"]),
    ]
    cw_each = CW / 4 - 0.25
    for i, (title, items) in enumerate(sectors):
        colors = [NAVY, BLUE, DARK2, DARK_GOLD]
        bgs    = [GREY, GREY, GREY, CREAM]
        card(s, ML + i * (cw_each + 0.3), y, cw_each, CH - 0.5,
             title, items, header_color=colors[i], body_color=bgs[i])


def slide_s01_problem(prs):
    s = content_slide(prs, "01", "El problema",
                      "El proceso actual: lento, manual y propenso a errores",
                      "Localizar todos los activos de una empresa grande requiere 2–4 horas de trabajo manual con cobertura parcial y sin trazabilidad del proceso.")
    y = CY + 0.3
    # Timeline 4 steps + pain points
    steps = [
        ("BÚSQUEDA", "Web, LinkedIn,\nRegistro Mercantil"),
        ("CONTRASTE", "Cruzar fuentes:\nBBDD, informes anuales"),
        ("GEOCODIFICACIÓN", "Buscar coordenadas\npara cada dirección"),
        ("VALIDACIÓN", "Verificar si el activo\nsigue activo"),
    ]
    sw_step = CW / 4 - 0.3
    for i, (title, desc) in enumerate(steps):
        x = ML + i * (sw_step + 0.35)
        # Step box
        s_step = _rect(s, x, y, sw_step, 3.5, fill=NAVY, rounding=True)
        _shape_txt(s_step, [title, "", desc], 10, color=WHITE, align=PP_ALIGN.CENTER)
        tf = s_step.text_frame
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.size = Pt(11)
        # Arrow between steps
        if i < len(steps) - 1:
            arrow_right(s, x + sw_step + 0.05, y + 1.3)

    # Pain points section
    _hline(s, ML, y + 4.0, CW, DARK_GOLD, 2)
    _txt(s, ML, y + 4.3, CW, 0.7, "⚠  Principales problemas del proceso manual", 11, bold=True, color=DARK_GOLD)
    pains = [
        ("Tiempo", "2–4 h por empresa para cobertura parcial"),
        ("Exactitud", "Activos cerrados, coordenadas incorrectas o ausentes"),
        ("Escalabilidad", "Imposible aplicar a portfolios de +100 empresas"),
        ("Trazabilidad", "Sin log de fuentes consultadas ni fecha de validación"),
    ]
    pw = CW / 4 - 0.3
    for i, (lbl, desc) in enumerate(pains):
        px = ML + i * (pw + 0.35)
        pb = _rect(s, px, y + 5.1, pw, 2.8, fill=CREAM, border=DARK_GOLD, border_w_pt=0.8, rounding=True)
        _shape_txt(pb, [lbl, "", desc], 9.5, color=BLACK, align=PP_ALIGN.CENTER)
        tf = pb.text_frame
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.color.rgb = RED_C


def slide_s01_limitations(prs):
    s = content_slide(prs, "01", "El problema",
                      "Limitaciones de las soluciones existentes",
                      "Ninguna herramienta actual combina la automatización, el filtrado por activos productivos y un scoring de confianza reproducible.")
    y = CY + 0.3
    two_col(s, y,
            "Fuentes y herramientas actuales",
            ["Google Maps: no distingue activos productivos de otros locales",
             "Registro Mercantil / CNAE: sin coordenadas ni verificación de estado",
             "LinkedIn: parcial, no estructurado, depende de actualizaciones del usuario",
             "Bureau van Dijk Orbis: costoso, sin API flexible, datos no geolocalizados",
             "Scrapers ad-hoc: frágiles, sin mantenimiento, sin scoring"],
            "Por qué no resuelven el problema",
            ["No filtran activos productivos vs. otros establecimientos",
             "No asignan un score de confianza por activo verificado",
             "No combinan múltiples fuentes heterogéneas automáticamente",
             "No ofrecen una UI interactiva de revisión y exportación",
             "No escalan a portfolios de cientos de empresas"],
            left_color=NAVY, right_color=RED_C, h=9.0)


def slide_s01_opportunity(prs):
    s = content_slide(prs, "01", "El problema",
                      "La oportunidad: convergencia de tres tecnologías clave",
                      "Por primera vez, APIs geoespaciales, modelos de lenguaje y agentes autónomos permiten automatizar completamente el proceso.")
    y = CY + 0.5
    pillars = [
        ("Google Maps\nPlaces API", "Cobertura global de establecimientos con datos estructurados: nombre, tipo, coordenadas, web, reseñas", GEO_BLUE),
        ("AWS Bedrock\n(Claude 3)", "Modelos de lenguaje capaces de clasificar, filtrar y valorar activos con razonamiento contextual a bajo coste", ORG_C),
        ("CrewAI +\nDuckDuckGo MCP", "Agentes autónomos que navegan la web, descargan y evalúan documentos sin intervención humana", TEAL_C),
    ]
    pw = CW / 3 - 0.4
    for i, (title, desc, color) in enumerate(pillars):
        px = ML + i * (pw + 0.55)
        # Large pillar block
        p_rect = _rect(s, px, y, pw, 6.0, fill=color, rounding=True)
        _shape_txt(p_rect, [title, "", desc], 10, color=WHITE, align=PP_ALIGN.CENTER)
        tf = p_rect.text_frame
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.size = Pt(13)

    # Central connector: "GeoAssets Intelligence"
    center_box = _rect(s, ML + CW/2 - 5.5, y + 6.5, 11, 1.4, fill=NAVY, rounding=True)
    _shape_txt(center_box, ["GeoAssets Intelligence"], 14, bold=True,
               color=WHITE, align=PP_ALIGN.CENTER)
    # Small arrows pointing down
    for i in range(3):
        px = ML + i * (pw + 0.55) + pw / 2 - 0.3
        arrow_right_s = s.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, Cm(px), Cm(y + 6.0), Cm(0.6), Cm(0.7))
        arrow_right_s.fill.solid(); arrow_right_s.fill.fore_color.rgb = GOLD
        arrow_right_s.line.fill.background()


def slide_s01_objective(prs):
    s = content_slide(prs, "01", "El problema",
                      "Objetivo y alcance del TFM",
                      "Construir un sistema end-to-end que localice, clasifique y puntúe los activos productivos de una empresa con mínima intervención humana.")
    y = CY + 0.3
    # Objective box
    obj_box = _rect(s, ML, y, CW, 2.2, fill=NAVY, rounding=True)
    _shape_txt(obj_box, ["OBJETIVO PRINCIPAL", "",
                         "Sistema end-to-end de localización automatizada de activos productivos empresariales con scoring de confianza reproducible"],
               11, color=WHITE, align=PP_ALIGN.CENTER)
    tf = obj_box.text_frame
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.size = Pt(12)

    # Scope 3 columns
    scopes = [
        ("Ámbito geográfico", ["España (expansión futura a EMEA)", "Google Maps + fuentes web en español", "12 categorías de activo productivo"]),
        ("KPIs de éxito", ["Cobertura ≥ 80 % de activos conocidos", "Tier HIGH ≥ 60 % de resultados", "Latencia < 3 min por empresa (Maps)"]),
        ("Fuera de alcance", ["Activos residenciales / inmuebles", "Empresas fuera de España (v1)", "Validación sobre el terreno"]),
    ]
    sw_s = CW / 3 - 0.35
    colors_s = [DARK_GOLD, GRN_C, MID_GREY]
    for i, (title, items) in enumerate(scopes):
        card(s, ML + i * (sw_s + 0.5), y + 2.6, sw_s, CH - 2.8,
             title, items, header_color=colors_s[i])


def slide_s02_platform(prs):
    s = content_slide(prs, "02", "La solución",
                      "GeoAssets Intelligence: presentación de la plataforma",
                      "Una plataforma web que transforma el nombre de una empresa en un mapa georreferenciado de sus activos productivos en minutos.")
    y = CY + 0.3
    # Left: description + tags
    desc_box = _rect(s, ML, y, 14.5, 4.0, fill=NAVY, rounding=True)
    _shape_txt(desc_box, [
        "Plataforma de inteligencia geoespacial", "",
        "Localiza · Clasifica · Puntúa", "",
        "activos productivos mediante tres pipelines de análisis basados en IA"
    ], 11, color=WHITE, align=PP_ALIGN.CENTER)
    tf = desc_box.text_frame
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.size = Pt(13)
    tf.paragraphs[2].runs[0].font.bold = True
    tf.paragraphs[2].runs[0].font.color.rgb = GOLD

    # Right: 3 highlight chips
    highlights = [
        ("⚡ Tiempo", "~2 minutos por empresa\n(modo Maps API)", GEO_BLUE),
        ("🎯 Precisión", "Score de confianza\npor señales objetivas", GRN_C),
        ("📤 Exportación", "CSV y Excel\ncon todos los campos", DARK_GOLD),
    ]
    for i, (lbl, desc, clr) in enumerate(highlights):
        bx = ML + 15.0
        by = y + i * 1.5
        b = _rect(s, bx, by, CW - 15.0, 1.35, fill=clr, rounding=True)
        _shape_txt(b, [lbl, desc], 10, color=WHITE, align=PP_ALIGN.LEFT)
        tf = b.text_frame
        tf.paragraphs[0].runs[0].font.bold = True

    # Screenshot placeholder area
    ph = _rect(s, ML, y + 4.3, CW, CH - 4.5, fill=GREY, border=NAVY, border_w_pt=1.0)
    _shape_txt(ph, [
        "[ Captura de pantalla del dashboard principal ]",
        "",
        "Mapa interactivo Leaflet  ·  Sidebar con lista y filtros  ·  Popup de activo con scoring"
    ], 11, color=MID_GREY, align=PP_ALIGN.CENTER)


def slide_s02_modes(prs):
    s = content_slide(prs, "02", "La solución",
                      "Tres modos de análisis para cualquier escenario",
                      "El usuario elige el modo según los datos disponibles: desde nombre de empresa hasta documentos internos o búsqueda web autónoma.")
    y = CY + 0.3
    modes = [
        ("01", "Búsqueda Maps API", GEO_BLUE,
         ["Input: solo el nombre de la empresa",
          "Pipeline automático de 5 pasos",
          "Google Places API + filtro LLM",
          "~45 segundos de media",
          "Ideal: empresas con presencia en Google Maps"]),
        ("02", "Carga de Documento", PURPLE_C,
         ["Input: PDF, DOCX o Excel con activos",
          "Pipeline de 6 pasos con parsing IA",
          "Extracción + geocodificación + scoring",
          "~60 segundos de media",
          "Ideal: memoria anual, informes internos"]),
        ("03", "Agente IA Autónomo", TEAL_C,
         ["Input: solo el nombre de la empresa",
          "Agente CrewAI navega la web",
          "Descarga y evalúa documentos relevantes",
          "~3 minutos de media",
          "Ideal: empresas sin datos en Google Maps"]),
    ]
    cw_m = CW / 3 - 0.35
    for i, (num, title, color, items) in enumerate(modes):
        mx = ML + i * (cw_m + 0.5)
        # Mode header
        hdr = _rect(s, mx, y, cw_m, 1.6, fill=color, rounding=True)
        _shape_txt(hdr, [f"MODO {num}", title], 11, color=WHITE, align=PP_ALIGN.CENTER)
        tf = hdr.text_frame
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.size = Pt(13)
        # Items
        body = _rect(s, mx, y + 1.6, cw_m, CH - 1.8, fill=GREY, border=color, border_w_pt=0.8)
        _shape_txt(body, items, 10, color=BLACK, align=PP_ALIGN.LEFT)
        from pptx.oxml.ns import qn
        txBody = body.text_frame._txBody
        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('lIns', str(int(Cm(0.3))))
            bodyPr.set('tIns', str(int(Cm(0.3))))


def slide_s02_taxonomy(prs):
    s = content_slide(prs, "02", "La solución",
                      "Taxonomía de activos: 12 categorías predefinidas",
                      "El modelo LLM asigna a cada activo una de las 12 categorías del enum AssetCategory, con icon y color propios en el mapa.")
    y = CY + 0.3
    cats = [
        ("HQ", "Sede Central", NAVY),
        ("OFF", "Oficina Regional", GEO_BLUE),
        ("FAB", "Fábrica / Planta", RGBColor(0xE6, 0x51, 0x00)),
        ("LOG", "Centro Logístico", ORG_C),
        ("TEC", "Centro Tecnológico", PURPLE_C),
        ("COM", "Punto de Venta", GRN_C),
        ("AGR", "Explotación Agrícola", RGBColor(0x33, 0x69, 0x1E)),
        ("ENE", "Inst. Energética", RGBColor(0xFD, 0xD8, 0x35)),
        ("TRA", "Infraestructura Transporte", MID_GREY),
        ("HOT", "Activo Hotelero", RGBColor(0xE9, 0x1E, 0x63)),
        ("SAN", "Centro Sanitario", RED_C),
        ("OTR", "Otro", MID_GREY),
    ]
    cols = 4
    rows = 3
    cw_c = CW / cols - 0.3
    ch_c = (CH) / rows - 0.3
    for i, (code, label, color) in enumerate(cats):
        col = i % cols
        row = i // cols
        cx = ML + col * (cw_c + 0.35)
        cy = y + row * (ch_c + 0.35)
        b = _rect(s, cx, cy, cw_c, ch_c, fill=color, rounding=True)
        _shape_txt(b, [code, label], 10, color=WHITE, align=PP_ALIGN.CENTER)
        tf = b.text_frame
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.size = Pt(12)


def slide_s02_scoring(prs):
    s = content_slide(prs, "02", "La solución",
                      "Sistema de confianza: scoring con distribución Beta",
                      "Cada activo recibe un score 0–100 % calculado a partir de señales ponderadas, suavizado con distribución Beta para evitar extremos artificiales.")
    y = CY + 0.3
    # Tier chips row
    tiers = [
        ("HIGH  ≥ 70 %", GRN_C, "Activo claramente perteneciente a la empresa"),
        ("MEDIUM  40–69 %", ORG_C, "Activo probable pero con señales mixtas"),
        ("LOW  < 40 %", RED_C, "Activo dudoso, requiere revisión manual"),
    ]
    tw = CW / 3 - 0.4
    for i, (lbl, clr, desc) in enumerate(tiers):
        tx = ML + i * (tw + 0.55)
        b = _rect(s, tx, y, tw, 1.0, fill=clr, rounding=True)
        _shape_txt(b, [lbl], 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(s, tx, y + 1.1, tw, 0.8, desc, 9.5, italic=True, color=BLACK)

    # Two signal tables side by side
    y2 = y + 2.3
    maps_signals = [
        ("Nombre incluye la empresa", "30 %"),
        ("Tipo compatible con activo productivo", "20 %"),
        ("Señal de ubicación corporativa", "15 %"),
        ("Dominio web corporativo", "15 %"),
        ("Perfil de reseñas B2B", "10 %"),
        ("Valoración del modelo IA", "10 %"),
    ]
    doc_signals = [
        ("Fuerza de la evidencia documental", "30 %"),
        ("Especificidad de la dirección", "20 %"),
        ("Origen de las coordenadas", "20 %"),
        ("Calidad del nombre del activo", "15 %"),
        ("Valoración del modelo IA", "15 %"),
    ]
    col_w2 = CW / 2 - 0.5

    # Maps table
    _rect(s, ML, y2, col_w2, 0.75, fill=GEO_BLUE)
    _txt(s, ML + 0.2, y2 + 0.1, col_w2 - 0.3, 0.6, "Pipeline Maps API – 6 señales", 10, bold=True, color=WHITE)
    for i, (sig, wt) in enumerate(maps_signals):
        row_color = GREY if i % 2 == 0 else WHITE
        row_y = y2 + 0.75 + i * 0.75
        _rect(s, ML, row_y, col_w2, 0.75, fill=row_color, border=GREY, border_w_pt=0.3)
        _txt(s, ML + 0.2, row_y + 0.1, col_w2 - 2.5, 0.6, sig, 9.5, color=BLACK)
        b = _rect(s, ML + col_w2 - 2.0, row_y + 0.1, 1.8, 0.6, fill=GEO_BLUE, rounding=True)
        _shape_txt(b, [wt], 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Doc table
    rx = ML + col_w2 + 0.5
    _rect(s, rx, y2, col_w2, 0.75, fill=PURPLE_C)
    _txt(s, rx + 0.2, y2 + 0.1, col_w2 - 0.3, 0.6, "Pipeline Documento – 5 señales", 10, bold=True, color=WHITE)
    for i, (sig, wt) in enumerate(doc_signals):
        row_color = GREY if i % 2 == 0 else WHITE
        row_y = y2 + 0.75 + i * 0.75
        _rect(s, rx, row_y, col_w2, 0.75, fill=row_color, border=GREY, border_w_pt=0.3)
        _txt(s, rx + 0.2, row_y + 0.1, col_w2 - 2.5, 0.6, sig, 9.5, color=BLACK)
        b = _rect(s, rx + col_w2 - 2.0, row_y + 0.1, 1.8, 0.6, fill=PURPLE_C, rounding=True)
        _shape_txt(b, [wt], 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_s03_stack(prs):
    s = content_slide(prs, "03", "Arquitectura técnica",
                      "Stack tecnológico",
                      "Combinación de tecnologías modernas de código abierto con servicios cloud gestionados para máxima velocidad de desarrollo y escalabilidad.")
    y = CY + 0.3
    tech = [
        ("Frontend", [("Vue 3 + Vuetify 3", GEO_BLUE), ("Pinia (state management)", GEO_BLUE),
                      ("Leaflet.js (mapas interactivos)", GEO_BLUE)], NAVY),
        ("Backend", [("FastAPI (Python 3.12)", DARK2), ("Server-Sent Events (SSE)", DARK2),
                     ("LiteLLM (proxy LLM)", DARK2)], DARK2),
        ("IA & Agentes", [("AWS Bedrock – Claude 3 Sonnet/Haiku", ORG_C),
                           ("CrewAI – orquestación de agentes", TEAL_C),
                           ("DuckDuckGo MCP – búsqueda web", TEAL_C)], ORG_C),
        ("Datos & Infra", [("Redis – caché TTL 24h", GRN_C), ("PostgreSQL – persistencia", GRN_C),
                            ("Docker – contenedorización", MID_GREY)], GRN_C),
    ]
    cw_t = CW / 4 - 0.3
    for i, (layer, items, color) in enumerate(tech):
        tx = ML + i * (cw_t + 0.35)
        hdr = _rect(s, tx, y, cw_t, 0.85, fill=color, rounding=True)
        _shape_txt(hdr, [layer], 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for j, (item, clr) in enumerate(items):
            iy = y + 0.85 + j * 1.0 + 0.15
            b = _rect(s, tx, iy, cw_t, 0.85, fill=GREY, border=color, border_w_pt=0.5, rounding=True)
            _shape_txt(b, [item], 9.5, color=clr, align=PP_ALIGN.CENTER)

    # Bottom note
    _txt(s, ML, y + CH - 0.9, CW, 0.8,
         "✓  Sin vendor lock-in en LLM: LiteLLM permite cambiar de modelo sin modificar el código de negocio",
         10, italic=True, color=DARK_GOLD)


def slide_s03_architecture(prs):
    s = content_slide(prs, "03", "Arquitectura técnica",
                      "Diagrama de arquitectura: tres capas",
                      "Arquitectura en capas con flujo de datos de petición (→) y respuesta en streaming (SSE ←), desplegada en contenedores Docker.")
    y = CY + 0.3

    layers = [
        ("PRESENTACIÓN", "Vue 3 + Vuetify + Pinia + Leaflet", GEO_BLUE,
         ["EventSource (SSE)", "Pinia store reactivo", "Mapa Leaflet + Popups"]),
        ("API & LÓGICA", "FastAPI + Orquestadores de pipeline", NAVY,
         ["run_pipeline_sse()", "run_doc_pipeline_sse()", "SSE streaming a cliente"]),
        ("SERVICIOS", "IA + Datos + APIs externas", DARK_GOLD,
         ["AWS Bedrock (Claude)", "Google Maps / Places API", "Redis + PostgreSQL"]),
    ]
    layer_w = CW / 3 - 0.4
    for i, (lbl, subtitle, clr, items) in enumerate(layers):
        lx = ML + i * (layer_w + 0.55)
        # Layer header
        lhdr = _rect(s, lx, y, layer_w, 1.3, fill=clr, rounding=True)
        _shape_txt(lhdr, [lbl, subtitle], 10, color=WHITE, align=PP_ALIGN.CENTER)
        tf = lhdr.text_frame
        tf.paragraphs[0].runs[0].font.bold = True
        # Items
        for j, item in enumerate(items):
            iy = y + 1.5 + j * 1.2
            ib = _rect(s, lx, iy, layer_w, 1.05, fill=GREY, border=clr, border_w_pt=0.6, rounding=True)
            _shape_txt(ib, [item], 9.5, color=BLACK, align=PP_ALIGN.CENTER)
        # Arrow to next layer
        if i < len(layers) - 1:
            ax = lx + layer_w + 0.05
            a = s.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                Cm(ax), Cm(y + 0.5), Cm(0.45), Cm(0.6))
            a.fill.solid(); a.fill.fore_color.rgb = GOLD
            a.line.fill.background()
            # SSE arrow (reverse)
            a2 = s.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.LEFT_ARROW,
                Cm(ax), Cm(y + 1.4), Cm(0.45), Cm(0.5))
            a2.fill.solid(); a2.fill.fore_color.rgb = TEAL_C
            a2.line.fill.background()

    # SSE protocol
    y2 = y + 5.5
    _hline(s, ML, y2, CW, NAVY, 2)
    _txt(s, ML, y2 + 0.2, CW, 0.7, "Protocolo SSE: job_started → step_start → step_complete → complete / error", 10, bold=True, color=NAVY)
    _txt(s, ML, y2 + 1.0, CW, 0.7, "El frontend suscribe a EventSource; cada evento actualiza el estado de Pinia y la barra de progreso de la UI en tiempo real", 9.5, color=BLACK)


def slide_s03_data_model(prs):
    s = content_slide(prs, "03", "Arquitectura técnica",
                      "Modelo de datos y gestión de caché",
                      "Dos entidades principales (Company, Asset) persistidas en PostgreSQL con caché Redis TTL=24h para evitar llamadas redundantes a las APIs.")
    y = CY + 0.3
    # Asset entity fields
    asset_fields = [
        ("id, company_id", "UUIDs de activo y empresa", GREY),
        ("name, raw_name", "Nombre normalizado y nombre original extraído", GREY),
        ("category", "Enum AssetCategory (12 valores)", GREY),
        ("latitude, longitude", "Coordenadas WGS84", GREY),
        ("address, municipality, province", "Dirección completa estructurada", CREAM),
        ("confidence_score", "Float 0–1 (score Beta suavizado)", CREAM),
        ("confidence_tier", "HIGH / MEDIUM / LOW", CREAM),
        ("confidence_signals", "Dict[str, float] — señales individuales", CREAM),
        ("data_sources", "List[str] — pipelines que aportaron datos", CREAM),
        ("website, phone", "Datos de contacto (Maps enrichment)", GREY),
    ]
    fh = (CH - 1.4) / len(asset_fields)
    hdr = _rect(s, ML, y, CW * 0.6, 0.75, fill=NAVY)
    _txt(s, ML + 0.2, y + 0.1, CW * 0.6 - 0.3, 0.6, "Entidad Asset — campos principales", 10, bold=True, color=WHITE)
    for i, (field, desc, bg) in enumerate(asset_fields):
        ry = y + 0.75 + i * fh
        row = _rect(s, ML, ry, CW * 0.6, fh, fill=bg, border=GREY, border_w_pt=0.3)
        _txt(s, ML + 0.2, ry + 0.05, 8, fh - 0.1, field, 9, bold=True, color=NAVY)
        _txt(s, ML + 8.3, ry + 0.05, CW * 0.6 - 8.5, fh - 0.1, desc, 9, color=BLACK)

    # Cache strategy
    cx = ML + CW * 0.6 + 0.5
    cw2 = CW * 0.4 - 0.5
    cache_items = [
        ("Redis TTL 24h", "Evita llamadas repetidas a Maps API y Bedrock para la misma empresa", GRN_C),
        ("force_refresh=True", "El usuario puede forzar re-análisis desde la UI", ORG_C),
        ("PostgreSQL", "Persistencia histórica: series temporales de activos por empresa", GEO_BLUE),
    ]
    for i, (title, desc, clr) in enumerate(cache_items):
        iy = y + i * 3.5
        b = _rect(s, cx, iy, cw2, 3.2, fill=clr, rounding=True)
        _shape_txt(b, [title, "", desc], 10, color=WHITE, align=PP_ALIGN.CENTER)
        tf = b.text_frame
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.size = Pt(11)


def slide_s04_overview(prs):
    s = content_slide(prs, "04", "Pipelines de análisis",
                      "Visión general: tres pipelines complementarios",
                      "Cada pipeline está optimizado para un tipo de datos de entrada diferente; el modo combinado puede fusionar resultados de varios pipelines.")
    y = CY + 0.2
    headers = ["", "Maps API", "Documento", "Agente IA"]
    rows = [
        ("Fuente de datos", "Google Places API", "PDF / DOCX / Excel", "Búsqueda web DuckDuckGo"),
        ("Nº de pasos", "5 pasos", "6 pasos", "Variable (agentic)"),
        ("Tiempo estimado", "~45 segundos", "~60 segundos", "~3 minutos"),
        ("Precisión esperada", "Alta (tier HIGH ~65 %)", "Muy alta (tier HIGH ~75 %)", "Media-alta"),
        ("Caso de uso ideal", "Empresa en Google Maps", "Informe anual / BD interna", "Sin datos previos"),
        ("Requisito", "Nombre de empresa", "Fichero subido por usuario", "Solo nombre de empresa"),
    ]
    col_widths = [6.5, 7.5, 7.5, 7.5]
    col_colors = [NAVY, GEO_BLUE, PURPLE_C, TEAL_C]
    total_w = sum(col_widths) + 0.3 * 3
    start_x = ML + (CW - total_w) / 2

    # Header row
    cx = start_x
    for i, (hdr, wid, clr) in enumerate(zip(headers, col_widths, col_colors)):
        if i == 0:
            _rect(s, cx, y, wid, 0.85, fill=GREY)
        else:
            hb = _rect(s, cx, y, wid, 0.85, fill=clr)
            _shape_txt(hb, [hdr], 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cx += wid + 0.3

    # Data rows
    for ri, row in enumerate(rows):
        ry = y + 0.85 + ri * 1.05
        row_bg = GREY if ri % 2 == 0 else WHITE
        cx = start_x
        for ci, (cell, wid, clr) in enumerate(zip(row, col_widths, col_colors)):
            rb = _rect(s, cx, ry, wid, 1.05, fill=row_bg, border=GREY, border_w_pt=0.3)
            color_txt = clr if ci > 0 and ri == 0 else (NAVY if ci == 0 else BLACK)
            bold_txt = ci == 0
            _txt(s, cx + 0.15, ry + 0.1, wid - 0.25, 0.9, cell, 9.5,
                 bold=bold_txt, color=color_txt)
            cx += wid + 0.3


def slide_s04_maps_pipeline(prs):
    s = content_slide(prs, "04", "Pipelines de análisis",
                      "Pipeline Maps API: 5 pasos automatizados",
                      "El pipeline más rápido: ingiere solo el nombre de la empresa y retorna activos geocodificados con scoring en ~45 segundos.")
    y = CY + 0.5
    steps = [
        ("0", "Identificar\nempresa", "LLM extrae nombre\ncanónico, CIF, sector,\nCNAE y sede"),
        ("1", "Buscar en\nGoogle Maps", "Múltiples queries con\nvariantes del nombre;\ndevuelve hasta N places"),
        ("2", "Filtrar y\nClasificar", "LLM descarta locales\nno productivos y asigna\ncategoría del enum"),
        ("3", "Enriquecer\ndatos", "Places Details API:\nweb, teléfono, reseñas\ny señales B2B"),
        ("4", "Calcular\nconfianza", "Score ponderado de 6\nseñales + suavizado\nBeta → tier HIGH/MED/LOW"),
    ]
    sw_s = CW / 5 - 0.5
    for i, (num, title, desc) in enumerate(steps):
        x = ML + i * (sw_s + 0.6)
        pipeline_step(s, x, y, sw_s, CH - 0.5, num, title, desc)
        if i < len(steps) - 1:
            arrow_right(s, x + sw_s + 0.05, y + CH / 2 - 0.5)

    # Confidence signals row at bottom
    y2 = y + CH + 0.2
    _txt(s, ML, y2 - 0.5, CW, 0.6,
         "Señales del score: name_match (30%) · type_match (20%) · address_corporate (15%) · website_match (15%) · reviews_b2b (10%) · llm_confidence (10%)",
         8.5, italic=True, color=DARK_GOLD)


def slide_s04_doc_pipeline(prs):
    s = content_slide(prs, "04", "Pipelines de análisis",
                      "Pipeline Documento: extracción IA de activos de ficheros",
                      "El pipeline de mayor precisión: procesa informes anuales, memorias de activos o bases de datos internas para extraer localizaciones con alta fiabilidad.")
    y = CY + 0.5
    steps = [
        ("0", "Parsear\ndocumento", "PDF/DOCX/Excel\n→ Markdown\n(LlamaParse/pdfminer)"),
        ("1", "Chunking\nsemántico", "Divide Markdown\nen chunks\ncoherentes"),
        ("2", "Extraer\nactivos (IA)", "LLM extrae entidades\nnombre, dirección, tipo\nde cada chunk"),
        ("3", "Deduplicar", "Agrupa menciones\ndel mismo activo por\nsimilitud de nombre"),
        ("4", "Geocodificar\ny enriquecer", "Google Geocoding API\n+ Places Details\npara coordenadas"),
        ("5", "Calcular\nconfianza", "5 señales documentales\n+ suavizado Beta\n→ tier HIGH/MED/LOW"),
    ]
    sw_s = CW / 6 - 0.4
    for i, (num, title, desc) in enumerate(steps):
        x = ML + i * (sw_s + 0.5)
        pipeline_step(s, x, y, sw_s, CH - 0.5, num, title, desc, color=PURPLE_C)
        if i < len(steps) - 1:
            arrow_right(s, x + sw_s + 0.05, y + CH / 2 - 0.5)

    _txt(s, ML, y + CH - 0.3, CW, 0.55,
         "Señales: evidence_strength (30%) · address_specificity (20%) · coordinate_source (20%) · name_quality (15%) · llm_confidence (15%)",
         8.5, italic=True, color=PURPLE_C)


def slide_s04_agent(prs):
    s = content_slide(prs, "04", "Pipelines de análisis",
                      "Pipeline Agente IA: búsqueda web autónoma",
                      "El agente actúa como un analista humano: lanza búsquedas, evalúa documentos y los filtra por relevancia antes de pasarlos al pipeline de documento.")
    y = CY + 0.3
    # Two phases
    ph1_items = [
        "1. AgentSearchView: el agente lanza queries en DuckDuckGo",
        "2. Descarga PDFs y DOCXs de URLs encontradas",
        "3. Emite eventos en tiempo real: thinking → searching → found_urls",
        "4. Acepta o rechaza cada documento por criterio de relevancia",
        "5. El usuario revisa la lista antes de confirmar el análisis",
    ]
    ph2_items = [
        "1. AgentDocumentReviewView: lista de ficheros aceptados",
        "2. El usuario confirma → lanza pipeline de documento por cada fichero",
        "3. Procesamiento paralelo configurable (max_concurrent)",
        "4. Fusión de resultados de múltiples fuentes en un único mapa",
        "5. Score de confianza incluye data_source = 'agent_search'",
    ]
    two_col(s, y,
            "Fase 1 – Búsqueda y selección de documentos",
            ph1_items,
            "Fase 2 – Revisión y análisis de documentos",
            ph2_items,
            left_color=TEAL_C, right_color=DARK_GOLD, h=9.5)

    # Event stream chip row
    events = ["thinking", "searching", "found_urls", "downloading", "accepted", "rejected"]
    event_colors = [DARK2, GEO_BLUE, NAVY, DARK_GOLD, GRN_C, RED_C]
    ew = CW / len(events) - 0.35
    ey = y + 10.0
    _txt(s, ML, ey - 0.4, CW, 0.4, "Tipos de evento SSE del agente:", 9, bold=True, color=NAVY)
    for i, (ev, clr) in enumerate(zip(events, event_colors)):
        b = _rect(s, ML + i * (ew + 0.35), ey, ew, 0.65, fill=clr, rounding=True)
        _shape_txt(b, [ev], 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_s04_decision(prs):
    s = content_slide(prs, "04", "Pipelines de análisis",
                      "¿Qué pipeline utilizar? Árbol de decisión",
                      "La elección del pipeline depende de los datos disponibles y del tiempo aceptable; el modo combinado permite fusionar resultados de varios pipelines.")
    y = CY + 0.4
    # Decision boxes
    q1 = _rect(s, ML + CW/2 - 7, y, 14, 1.2, fill=NAVY, rounding=True)
    _shape_txt(q1, ["¿Tienes un documento con los activos? (PDF, DOCX, Excel)"],
               11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Yes branch (left)
    a1_y = _rect(s, ML + 2, y + 2.5, 5.5, 1.0, fill=GRN_C, rounding=True)
    _shape_txt(a1_y, ["SÍ"], 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    doc_box = _rect(s, ML + 0.5, y + 4.0, 8.5, 2.8, fill=PURPLE_C, rounding=True)
    _shape_txt(doc_box, ["PIPELINE DOCUMENTO", "", "Alta precisión · ~60 s\nFuente: documento interno"],
               11, color=WHITE, align=PP_ALIGN.CENTER)
    tf = doc_box.text_frame; tf.paragraphs[0].runs[0].font.bold = True

    # No branch middle
    q2 = _rect(s, ML + CW/2 - 7, y + 2.5, 14, 1.2, fill=DARK2, rounding=True)
    _shape_txt(q2, ["¿La empresa tiene presencia en Google Maps?"],
               11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Yes Maps
    maps_box = _rect(s, ML + CW/2 - 4.5, y + 5.5, 9, 2.8, fill=GEO_BLUE, rounding=True)
    _shape_txt(maps_box, ["PIPELINE MAPS API", "", "Rápido · ~45 s\nFuente: Google Places"],
               11, color=WHITE, align=PP_ALIGN.CENTER)
    tf = maps_box.text_frame; tf.paragraphs[0].runs[0].font.bold = True

    # No → Agent
    agent_box = _rect(s, ML + CW - 9.5, y + 5.5, 9, 2.8, fill=TEAL_C, rounding=True)
    _shape_txt(agent_box, ["PIPELINE AGENTE IA", "", "Autónomo · ~3 min\nFuente: búsqueda web"],
               11, color=WHITE, align=PP_ALIGN.CENTER)
    tf = agent_box.text_frame; tf.paragraphs[0].runs[0].font.bold = True

    # Combined mode note
    combined = _rect(s, ML, y + 9.5, CW, 1.5, fill=GOLD, rounding=True)
    _shape_txt(combined, ["MODO COMBINADO: es posible fusionar los resultados de varios pipelines en un único mapa con deduplicación"],
               11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def slide_s05_demo_inditex(prs):
    s = content_slide(prs, "05", "Demostración y resultados",
                      "Caso de prueba 1: Inditex — Pipeline Maps API",
                      "Empresa con alta presencia en Google Maps. El pipeline identifica la sede de A Coruña, centros logísticos y excluye tiendas minoristas franquiciadas.")
    y = CY + 0.3
    results = [
        ("Activos identificados", "18 activos", GEO_BLUE),
        ("Tier HIGH", "12 (67 %)", GRN_C),
        ("Tier MEDIUM", "4 (22 %)", ORG_C),
        ("Tier LOW", "2 (11 %)", RED_C),
        ("Tiempo de ejecución", "~42 segundos", DARK2),
    ]
    rw = CW / len(results) - 0.3
    for i, (lbl, val, clr) in enumerate(results):
        rx = ML + i * (rw + 0.35)
        b = _rect(s, rx, y, rw, 2.2, fill=clr, rounding=True)
        _shape_txt(b, [val, "", lbl], 10, color=WHITE, align=PP_ALIGN.CENTER)
        tf = b.text_frame
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.size = Pt(18)

    # Assets found
    assets = [
        ("Sede Central A Coruña", "HQ", GRN_C, "97 %"),
        ("Inditex Campus, Arteixo", "LOG", GRN_C, "91 %"),
        ("Plataforma Logística Zaragoza", "LOG", GRN_C, "88 %"),
        ("Centro Tecnológico Narón", "TEC", ORG_C, "72 %"),
        ("Almacén Meco (Madrid)", "LOG", ORG_C, "65 %"),
        ("Oficinas Barcelona", "OFF", ORG_C, "58 %"),
    ]
    ah = (CH - 2.8) / len(assets)
    _rect(s, ML, y + 2.5, CW, 0.7, fill=NAVY)
    _txt(s, ML + 0.2, y + 2.6, 14, 0.55, "Activo identificado", 9, bold=True, color=WHITE)
    _txt(s, ML + 14.5, y + 2.6, 4, 0.55, "Categoría", 9, bold=True, color=WHITE)
    _txt(s, ML + 19, y + 2.6, 5, 0.55, "Confianza", 9, bold=True, color=WHITE)
    _txt(s, ML + 24.5, y + 2.6, CW - 25, y + 2.6, "Tier", 9, bold=True, color=WHITE)

    for i, (name, cat, clr, score) in enumerate(assets):
        ry = y + 3.2 + i * ah
        bg = GREY if i % 2 == 0 else WHITE
        _rect(s, ML, ry, CW, ah, fill=bg, border=GREY, border_w_pt=0.3)
        _txt(s, ML + 0.2, ry + 0.05, 14, ah - 0.1, name, 9.5, color=BLUE)
        cat_b = _rect(s, ML + 14.5, ry + 0.05, 3.5, ah - 0.15, fill=NAVY, rounding=True)
        _shape_txt(cat_b, [cat], 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        score_b = _rect(s, ML + 19, ry + 0.05, 4, ah - 0.15, fill=clr, rounding=True)
        _shape_txt(score_b, [score], 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tier_txt = "HIGH" if clr == GRN_C else "MEDIUM"
        _txt(s, ML + 24.0, ry + 0.05, 6, ah - 0.1, tier_txt, 9, bold=True, color=clr)


def slide_s05_metrics(prs):
    s = content_slide(prs, "05", "Demostración y resultados",
                      "Métricas de calidad y rendimiento del sistema",
                      "Evaluación cuantitativa sobre un conjunto de empresas piloto con activos conocidos previamente validados de forma manual.")
    y = CY + 0.3
    metrics = [
        ("Precision@K", "82 %", "% activos HIGH-tier validados manualmente como correctos", GRN_C),
        ("Recall estimado", "78 %", "% activos conocidos recuperados vs. total real", GEO_BLUE),
        ("Latencia Maps", "~45 s", "Mediana de tiempo de ejecución del pipeline Maps", DARK2),
        ("Latencia Documento", "~60 s", "Mediana de tiempo (fichero 20 páginas)", PURPLE_C),
        ("Latencia Agente", "~180 s", "Mediana incluyendo búsqueda web", TEAL_C),
        ("Tier HIGH rate", "63 %", "Proporción de activos con confianza alta", DARK_GOLD),
    ]
    mw = CW / 3 - 0.4
    mh = 3.5
    for i, (name, value, desc, clr) in enumerate(metrics):
        col = i % 3
        row = i // 3
        mx = ML + col * (mw + 0.55)
        my = y + row * (mh + 0.4)
        b = _rect(s, mx, my, mw, mh, fill=GREY, border=clr, border_w_pt=1.5, rounding=True)
        _txt(s, mx + 0.3, my + 0.3, mw - 0.5, 1.5, value, 28, bold=True, color=clr)
        _txt(s, mx + 0.3, my + 1.6, mw - 0.5, 0.65, name, 10, bold=True, color=NAVY)
        _txt(s, mx + 0.3, my + 2.3, mw - 0.5, 1.0, desc, 9, italic=True, color=MID_GREY)


def slide_s05_confidence_ui(prs):
    s = content_slide(prs, "05", "Demostración y resultados",
                      "Interfaz de usuario: detalle del score de confianza",
                      "El popup del mapa muestra score, tier y un botón de información que abre un diálogo con el desglose completo de señales.")
    y = CY + 0.3
    # Popup mockup
    popup = _rect(s, ML, y, 13.5, CH, fill=WHITE, border=NAVY, border_w_pt=1.5, rounding=True)

    # Popup header
    _rect(s, ML, y, 13.5, 0.85, fill=NAVY, rounding=True)
    _txt(s, ML + 0.3, y + 0.1, 13, 0.65, "Plataforma Logística Zaragoza", 10, bold=True, color=WHITE)

    _rect(s, ML + 0.3, y + 1.0, 3, 0.6, fill=ORG_C, rounding=True)
    _shape_txt(_rect(s, ML + 0.3, y + 1.0, 3, 0.6, fill=ORG_C, rounding=True), ["LOG  Centro Logístico"], 8.5, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(s, ML + 0.3, y + 1.8, 13, 0.5, "📍 Pol. Ind. Plaza, Zaragoza", 8.5, color=MID_GREY)
    _txt(s, ML + 0.3, y + 2.4, 13, 0.5, "Centro de distribución con gestión automatizada de almacén", 8.5, italic=True, color=BLACK)
    _txt(s, ML + 0.3, y + 3.0, 4.5, 0.5, "Confianza:", 9, bold=True, color=BLACK)
    _rect(s, ML + 4.8, y + 3.05, 5.5, 0.4, fill=GRN_C, rounding=True)
    _txt(s, ML + 10.5, y + 3.0, 2.5, 0.5, "88 %", 9, bold=True, color=GRN_C)
    _txt(s, ML + 0.3, y + 3.7, 13, 0.5, "🏷  logística · distribución · almacén · ZARAGOZA", 8, color=DARK2)
    _rect(s, ML + 0.3, y + 4.5, 5.5, 0.55, fill=GEO_BLUE, rounding=True)
    _shape_txt(_rect(s, ML + 0.3, y + 4.5, 5.5, 0.55, fill=GEO_BLUE, rounding=True),
               ["Ver en Google Maps ↗"], 8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Divider → dialog
    arr = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        Cm(ML + 13.8), Cm(y + CH/2 - 0.4), Cm(1.2), Cm(0.8))
    arr.fill.solid(); arr.fill.fore_color.rgb = GOLD
    arr.line.fill.background()

    # Confidence dialog mockup
    dlg_x = ML + 15.5
    dialog = _rect(s, dlg_x, y, CW - 15.5, CH, fill=WHITE, border=NAVY, border_w_pt=1.5, rounding=True)
    _rect(s, dlg_x, y, CW - 15.5, 0.85, fill=NAVY, rounding=True)
    _txt(s, dlg_x + 0.3, y + 0.1, 14, 0.65, "Detalle de confianza", 10, bold=True, color=WHITE)

    score_bg = _rect(s, dlg_x + 0.3, y + 1.0, CW - 16.3, 1.5, fill=GREY, rounding=True)
    _txt(s, dlg_x + 0.5, y + 1.1, 5, 0.65, "88 %", 22, bold=True, color=GRN_C)
    _rect(s, dlg_x + 4.5, y + 1.4, 4, 0.65, fill=GRN_C, rounding=True)
    _shape_txt(_rect(s, dlg_x + 4.5, y + 1.4, 4, 0.65, fill=GRN_C, rounding=True),
               ["Alta confianza"], 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _txt(s, dlg_x + 0.3, y + 2.7, CW - 16, 0.45, "Fuente del análisis", 8.5, bold=True, color=MID_GREY)
    src_b = _rect(s, dlg_x + 0.3, y + 3.2, 4.5, 0.6, fill=GEO_BLUE, rounding=True)
    _shape_txt(src_b, ["🗺 Google Maps"], 8.5, color=WHITE, align=PP_ALIGN.CENTER)

    _txt(s, dlg_x + 0.3, y + 4.0, CW - 16, 0.45, "Señales utilizadas", 8.5, bold=True, color=MID_GREY)
    sigs = [
        ("Nombre incluye empresa", "30 %", 0.9, GRN_C),
        ("Tipo compatible", "20 %", 0.85, GRN_C),
        ("Ubicación corporativa", "15 %", 0.5, ORG_C),
        ("Dominio web", "15 %", 0.8, GRN_C),
    ]
    for i, (sig, wt, val, clr) in enumerate(sigs):
        sy = y + 4.55 + i * 1.0
        _txt(s, dlg_x + 0.3, sy, 7, 0.45, sig, 8, color=BLACK)
        bar_full = _rect(s, dlg_x + 0.3, sy + 0.45, 8, 0.3, fill=GREY)
        _rect(s, dlg_x + 0.3, sy + 0.45, 8 * val, 0.3, fill=clr)


def slide_s06_benchmarking(prs):
    s = content_slide(prs, "06", "Benchmarking competitivo",
                      "Panorama de soluciones existentes y posicionamiento",
                      "GeoAssets Intelligence se diferencia por combinar tres fuentes heterogéneas con scoring unificado, diseño específico para activos productivos y código abierto.")
    y = CY + 0.2
    competitors = [
        ("GeoAssets\nIntelligence", [("Automatización", "★★★★★"), ("Filtro productivos", "★★★★★"),
                                      ("Scoring propio", "★★★★★"), ("Coste", "★★★★★"), ("Integración API", "★★★★")], GEO_BLUE),
        ("Google Maps\nPlatform (raw)", [("Automatización", "★★★"), ("Filtro productivos", "★"),
                                          ("Scoring propio", "★"), ("Coste", "★★★★"), ("Integración API", "★★★★★")], MID_GREY),
        ("Bureau van Dijk\nOrbis", [("Automatización", "★★★"), ("Filtro productivos", "★★★★"),
                                     ("Scoring propio", "★★"), ("Coste", "★"), ("Integración API", "★★★")], DARK2),
        ("CARTO", [("Automatización", "★★★★"), ("Filtro productivos", "★★"), ("Scoring propio", "★★"),
                   ("Coste", "★★"), ("Integración API", "★★★★")], ORG_C),
        ("Scraper\nad-hoc", [("Automatización", "★★"), ("Filtro productivos", "★★★"),
                              ("Scoring propio", "★"), ("Coste", "★★★★★"), ("Integración API", "★")], RED_C),
    ]
    cw_c = CW / len(competitors) - 0.3
    criteria = ["Automatización", "Filtro productivos", "Scoring propio", "Coste", "Integración API"]
    row_h = (CH - 1.3) / (len(criteria) + 1)

    for i, (name, scores, clr) in enumerate(competitors):
        cx = ML + i * (cw_c + 0.35)
        # Header
        hb = _rect(s, cx, y, cw_c, 1.3, fill=clr, rounding=True)
        _shape_txt(hb, [name], 9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Score rows
        for j, (crit, stars) in enumerate(scores):
            ry = y + 1.3 + j * row_h
            bg = GREY if j % 2 == 0 else WHITE
            rb = _rect(s, cx, ry, cw_c, row_h, fill=bg, border=GREY, border_w_pt=0.3)
            star_color = GRN_C if clr == GEO_BLUE else MID_GREY
            _shape_txt(rb, [stars], 10, bold=True, color=star_color, align=PP_ALIGN.CENTER)

    # Criteria labels on the left (override with text after building)
    for j, crit in enumerate(criteria):
        ry = y + 1.3 + j * row_h + row_h/2 - 0.2
        # Left side label already in table cells — add floating label outside
        _txt(s, ML - 1.3, ry, 1.2, 0.4, crit, 7, italic=True, color=MID_GREY)


def slide_s07_value_prop(prs):
    s = content_slide(prs, "07", "Propuesta de valor",
                      "Propuesta de valor: del proceso manual al proceso automatizado",
                      "GeoAssets Intelligence transforma un proceso de 4 horas en 2 minutos con mayor cobertura, trazabilidad y scoring objetivo.")
    y = CY + 0.2
    two_col(s, y,
            "AS IS  —  Proceso actual",
            ["2–4 horas por empresa, cobertura parcial",
             "Sin trazabilidad de fuentes consultadas",
             "Sin score de confianza por activo",
             "Imposible escalar a +100 empresas",
             "Resultados no reproducibles"],
            "TO BE  —  Con GeoAssets Intelligence",
            ["~2 minutos por empresa, cobertura alta",
             "Log de fuentes y fecha de análisis por activo",
             "Score 0–100 % con señales ponderadas y tier",
             "Portfolio de cientos de empresas en batch",
             "Resultados reproducibles y exportables CSV/Excel"],
            left_color=RED_C, right_color=GRN_C, h=9.5)

    # ROI note
    roi = _rect(s, ML, y + 10.0, CW, 1.5, fill=GOLD, rounding=True)
    _shape_txt(roi, ["Impacto estimado: x120 reducción de tiempo · cobertura +40 % · trazabilidad 100 %"],
               11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def slide_s07_roadmap(prs):
    s = content_slide(prs, "07", "Propuesta de valor",
                      "Líneas de actuación futura",
                      "GeoAssets v1 valida el concepto en España; las siguientes fases expanden la cobertura geográfica, mejoran la precisión y abren la plataforma como API.")
    y = CY + 0.3
    phases = [
        ("Fase 1\nCorto plazo", ["Expansión a Portugal y Latinoamérica", "Adaptación de fuentes de datos locales",
                                   "Fine-tuning del clasificador con datos validados"], NAVY),
        ("Fase 2\nMedio plazo", ["API pública REST con autenticación", "Integración en plataformas de M&A y crédito",
                                  "Panel de administración multi-empresa"], GEO_BLUE),
        ("Fase 3\nLargo plazo", ["Módulo de detección de cambios (alertas)", "Dashboard ESG de exposición de activos",
                                  "Análisis de series temporales del patrimonio"], TEAL_C),
        ("Mejoras\nContinuas", ["Nuevos modelos LLM al lanzarse (plug-and-play)", "Nuevas categorías de activo por sector",
                                 "Feedback loop de validación humana al modelo"], DARK_GOLD),
    ]
    cw_p = CW / 4 - 0.35
    for i, (phase, items, color) in enumerate(phases):
        px = ML + i * (cw_p + 0.45)
        card(s, px, y, cw_p, CH - 0.5, phase, items, header_color=color)


def slide_s07_limitations(prs):
    s = content_slide(prs, "07", "Propuesta de valor",
                      "Limitaciones actuales y estrategias de mitigación",
                      "El sistema es funcional y preciso en su dominio pero presenta limitaciones conocidas con estrategias de mitigación identificadas.")
    y = CY + 0.3
    items = [
        ("Dependencia de Google Maps API",
         "Datos incompletos en zonas rurales o países con poca cobertura",
         "Complementar con Pipeline Agente para empresas sin presencia en Maps", RED_C),
        ("LLM no determinista",
         "Clasificaciones incorrectas de activos ambiguos en el 5–10 % de casos",
         "Revisar manualmente activos con tier LOW; validación por prompt engineering", ORG_C),
        ("Latencia del agente IA",
         "~3 minutos de espera puede ser inaceptable en algunos flujos",
         "Uso asíncrono con notificación al completar; pipeline Maps como alternativa rápida", ORG_C),
        ("Cobertura geográfica",
         "España únicamente en v1; otras geografías requieren adaptación de prompts y fuentes",
         "Arquitectura modular permite añadir fuentes por país sin refactoring mayor", DARK_GOLD),
    ]
    h_each = CH / len(items) - 0.3
    for i, (lim, problem, mitigation, clr) in enumerate(items):
        iy = y + i * (h_each + 0.3)
        # Limitation header
        _rect(s, ML, iy, CW, 0.7, fill=clr)
        _txt(s, ML + 0.3, iy + 0.1, CW - 0.5, 0.55, lim, 10, bold=True, color=WHITE)
        # Problem / Mitigation
        half_w = CW / 2 - 0.3
        pb = _rect(s, ML, iy + 0.7, half_w, h_each - 0.7, fill=CREAM, border=clr, border_w_pt=0.5)
        _shape_txt(pb, ["⚠ " + problem], 9.5, color=BLACK, align=PP_ALIGN.LEFT)
        from pptx.oxml.ns import qn
        for shape in [pb]:
            txBody = shape.text_frame._txBody
            bodyPr = txBody.find(qn('a:bodyPr'))
            if bodyPr is not None:
                bodyPr.set('lIns', str(int(Cm(0.2))))
                bodyPr.set('tIns', str(int(Cm(0.15))))
        mb = _rect(s, ML + half_w + 0.25, iy + 0.7, half_w, h_each - 0.7, fill=GREY, border=GRN_C, border_w_pt=0.5)
        _shape_txt(mb, ["✓ " + mitigation], 9.5, color=BLACK, align=PP_ALIGN.LEFT)
        for shape in [mb]:
            txBody = shape.text_frame._txBody
            bodyPr = txBody.find(qn('a:bodyPr'))
            if bodyPr is not None:
                bodyPr.set('lIns', str(int(Cm(0.2))))
                bodyPr.set('tIns', str(int(Cm(0.15))))


def slide_s07_conclusions(prs):
    s = content_slide(prs, "07", "Propuesta de valor",
                      "Conclusiones",
                      "GeoAssets Intelligence valida que los LLMs, en combinación con APIs geoespaciales y agentes autónomos, pueden automatizar el análisis de patrimonio empresarial.")
    y = CY + 0.3
    conclusions = [
        ("Validación de hipótesis", "Los LLMs + APIs geoespaciales permiten automatizar el análisis de activos productivos con alta precisión (Precision@K = 82 %)", GEO_BLUE),
        ("Contribución técnica", "Sistema modular de tres pipelines con scoring unificado basado en distribución Beta, UI interactiva y exportación de resultados", NAVY),
        ("Contribución metodológica", "Framework reproducible y extensible a otros dominios de inteligencia geoespacial empresarial", DARK_GOLD),
        ("Impacto práctico", "Reducción del tiempo de análisis de horas a minutos; aplicable a banca, seguros, consultoría estratégica y ESG", GRN_C),
    ]
    ch_each = (CH - 0.5) / len(conclusions)
    for i, (title, desc, clr) in enumerate(conclusions):
        iy = y + i * (ch_each + 0.15)
        # Check icon box
        icon_b = _rect(s, ML, iy, 1.5, ch_each - 0.1, fill=clr, rounding=True)
        _shape_txt(icon_b, ["✓"], 22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Content
        content_b = _rect(s, ML + 1.7, iy, CW - 1.7, ch_each - 0.1, fill=GREY, border=clr, border_w_pt=0.8, rounding=True)
        _shape_txt(content_b, [title, "", desc], 10, color=BLACK, align=PP_ALIGN.LEFT)
        tf = content_b.text_frame
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.color.rgb = clr
        from pptx.oxml.ns import qn
        txBody = content_b.text_frame._txBody
        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('lIns', str(int(Cm(0.3))))
            bodyPr.set('tIns', str(int(Cm(0.2))))


def slide_closing(prs):
    slide = add_slide(prs)
    add_bg(slide, NAVY)
    _rect(slide, 0, 0, 0.5, SH, fill=GOLD)
    _hline(slide, 0, SH / 2 - 0.1, SW, GOLD, 3)
    _hline(slide, 0, SH / 2 + 0.4, SW, DARK_GOLD, 1.5)

    _txt(slide, 2.0, 3.5, SW - 4, 3.5,
         "Gracias por vuestra atención", 30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(slide, 2.0, 7.0, SW - 4, 2.0,
         "GeoAssets Intelligence", 22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    _txt(slide, 2.0, 9.0, SW - 4, 1.5,
         "Localización automatizada del patrimonio productivo empresarial mediante IA",
         13, italic=True, color=DARK_GOLD, align=PP_ALIGN.CENTER)
    _hline(slide, 4, 11.0, SW - 8, DARK_GOLD, 1.5)
    _txt(slide, 2.0, 11.5, SW - 4, 1.0,
         "Alonso Salgueiro Castelo  ·  alonso.salgueiro.castelo@msspain.com",
         11, color=MID_GREY, align=PP_ALIGN.CENTER)
    _txt(slide, 2.0, 13.0, SW - 4, 1.5,
         "\"El patrimonio productivo de una empresa es visible.\nSolo hay que saber dónde mirar.\"",
         13, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
    _txt(slide, SW - 5, SH - 1.5, 4, 0.8,
         "Mayo 2026", 10, color=MID_GREY, align=PP_ALIGN.RIGHT)
    add_footer_line(slide)


def slide_bibliography(prs):
    s = content_slide(prs, "08", "Bibliografía",
                      "Bibliografía y referencias",
                      "Principales fuentes académicas, técnicas y de documentación oficial utilizadas en el desarrollo del TFM.")
    y = CY + 0.2
    refs = [
        "[1]  Google Developers. Google Places API documentation. developers.google.com/maps/documentation/places",
        "[2]  Amazon Web Services. Amazon Bedrock documentation. docs.aws.amazon.com/bedrock",
        "[3]  LiteLLM. Proxy for 100+ LLMs. litellm.vercel.app",
        "[4]  CrewAI. Multi-agent AI framework. crewai.com",
        "[5]  Brown, T. et al. (2020). Language Models are Few-Shot Learners. NeurIPS 2020.",
        "[6]  Reynolds, L. & McDonell, K. (2021). Prompt Programming for Large Language Models. CHI EA 2021.",
        "[7]  Gelman, A. & Hill, J. (2006). Data Analysis Using Regression and Multilevel/Hierarchical Models. CUP.",
        "[8]  Vue.js. The Progressive JavaScript Framework. vuejs.org",
        "[9]  FastAPI. Modern, fast web framework for building APIs with Python. fastapi.tiangolo.com",
        "[10] python-pptx. Python library for creating and updating PowerPoint files. python-pptx.readthedocs.io",
    ]
    rh = (CH - 0.2) / len(refs)
    for i, ref in enumerate(refs):
        ry = y + i * rh
        bg = GREY if i % 2 == 0 else WHITE
        b = _rect(s, ML, ry, CW, rh - 0.05, fill=bg, border=GREY, border_w_pt=0.2)
        _shape_txt(b, [ref], 9.5, color=BLUE, align=PP_ALIGN.LEFT)
        from pptx.oxml.ns import qn
        txBody = b.text_frame._txBody
        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('lIns', str(int(Cm(0.25))))
            bodyPr.set('tIns', str(int(Cm(0.1))))


# ════════════════════════════════════════════════════════════════
# MAIN BUILD
# ════════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width  = Cm(SW)
    prs.slide_height = Cm(SH)

    # 1 – Cover
    slide_cover(prs)

    # 2 – Index
    slide_index(prs)

    # ── SECCIÓN 01 ───────────────────────────────────────────────
    section_divider(prs, "01", "El problema",
                    "Opacidad del patrimonio productivo empresarial")
    slide_s01_context(prs)
    slide_s01_problem(prs)
    slide_s01_limitations(prs)
    slide_s01_opportunity(prs)
    slide_s01_objective(prs)

    # ── SECCIÓN 02 ───────────────────────────────────────────────
    section_divider(prs, "02", "La solución",
                    "Plataforma GeoAssets Intelligence")
    slide_s02_platform(prs)
    slide_s02_modes(prs)
    slide_s02_taxonomy(prs)
    slide_s02_scoring(prs)

    # ── SECCIÓN 03 ───────────────────────────────────────────────
    section_divider(prs, "03", "Arquitectura técnica",
                    "Stack, capas y modelo de datos")
    slide_s03_stack(prs)
    slide_s03_architecture(prs)
    slide_s03_data_model(prs)

    # ── SECCIÓN 04 ───────────────────────────────────────────────
    section_divider(prs, "04", "Pipelines de análisis",
                    "Tres pipelines complementarios para cualquier escenario")
    slide_s04_overview(prs)
    slide_s04_maps_pipeline(prs)
    slide_s04_doc_pipeline(prs)
    slide_s04_agent(prs)
    slide_s04_decision(prs)

    # ── SECCIÓN 05 ───────────────────────────────────────────────
    section_divider(prs, "05", "Demostración y resultados",
                    "Casos de prueba, métricas e interfaz de usuario")
    slide_s05_demo_inditex(prs)
    slide_s05_metrics(prs)
    slide_s05_confidence_ui(prs)

    # ── SECCIÓN 06 ───────────────────────────────────────────────
    section_divider(prs, "06", "Benchmarking competitivo",
                    "Posicionamiento frente a soluciones existentes")
    slide_s06_benchmarking(prs)

    # ── SECCIÓN 07 ───────────────────────────────────────────────
    section_divider(prs, "07", "Propuesta de valor",
                    "Impacto, roadmap y conclusiones")
    slide_s07_value_prop(prs)
    slide_s07_roadmap(prs)
    slide_s07_limitations(prs)
    slide_s07_conclusions(prs)

    # Closing + bibliography
    slide_closing(prs)
    slide_bibliography(prs)

    out = "/home/genmsadmin/Projects/geoasset-location/GeoAssets_Intelligence_TFM.pptx"
    prs.save(out)
    print(f"✓ Saved {len(prs.slides)} slides → {out}")


if __name__ == "__main__":
    build()
